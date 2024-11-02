import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from transformers import T5Tokenizer, T5ForConditionalGeneration, BlipProcessor, BlipForConditionalGeneration, pipeline
import torch
import base64
from PIL import Image
from docx import Document
from pptx import Presentation
import io
from langchain.schema import Document
import pdfplumber

# Model and tokenizer loading for text summarization
checkpoint = "LaMini-Flan-T5-248M"
tokenizer = T5Tokenizer.from_pretrained("MBZUAI/LaMini-Flan-T5-248M", legacy=False)
base_model = T5ForConditionalGeneration.from_pretrained("MBZUAI/LaMini-Flan-T5-248M")

# Summarization pipeline (with caching for faster reuse)
@st.cache_resource
def load_summarization_pipeline():
    return pipeline(
        'summarization',
        model=base_model,
        tokenizer=tokenizer,
        device=0 if torch.cuda.is_available() else -1,
        max_length=100,
        min_length=50
    )

# Image description model and processor
@st.cache_resource
def load_image_description_model():
    processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    return processor, model

# File processing functions for PDF, DOCX, PPTX
def file_preprocessing(file, file_type):
    pages = []

    if file_type == 'pdf':
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(Document(page_content=text))  # Wrap text in a Document object
    elif file_type == 'docx':
        doc = Document(file)
        pages.append(Document(page_content=doc.text))  # Wrap DOCX text in a Document object
    elif file_type == 'pptx':
        presentation = Presentation(file)
        for slide in presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text + "\n"  # Concatenate all text from the slide
            if slide_text.strip():  # Check if there's any text
                pages.append(Document(page_content=slide_text))  # Wrap slide text in a Document object

    if not pages:
        st.error("No text found in the uploaded file.")
        return ""

    # Now we can safely use the text_splitter
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=200, chunk_overlap=50)
    texts = text_splitter.split_documents(pages)  # This will now work correctly
    final_texts = ""
    for text in texts:
        final_texts += text.page_content
    return final_texts

# Summarization function
def llm_pipeline(file_content, file_type):
    summarizer = load_summarization_pipeline()
    input_text = file_preprocessing(file_content, file_type)
    if input_text:
        result = summarizer(input_text)
        return result[0]['summary_text'] if result else "No summary generated."
    return "No input text available for summarization."

# Image description function
def describe_image(image_file):
    processor, model = load_image_description_model()
    image = Image.open(image_file)
    inputs = processor(image, return_tensors="pt").to(model.device)
    outputs = model.generate(**inputs)
    description = processor.decode(outputs[0], skip_special_tokens=True)
    return description

# Function to display PDF content in Streamlit
@st.cache_data
def displayPDF(file):
    with open(file, "rb") as f:
        base64_pdf = base64.b64encode(f.read()).decode('utf-8')
    pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="600" type="application/pdf"></iframe>'
    st.markdown(pdf_display, unsafe_allow_html=True)

# Streamlit code
st.set_page_config(layout="wide")

def main():
    st.title("Document and Image Analysis App with Language Model")

    # Document summarization section
    st.header("Document Summarization")
    uploaded_file = st.file_uploader("Upload your document (PDF, DOCX, PPTX)", type=['pdf', 'docx', 'pptx'])

    if uploaded_file is not None:
        # Determine file type and process
        file_type = uploaded_file.type.split('/')[-1]

        # Save the uploaded file temporarily
        temp_file_path = f"temp_{uploaded_file.name}"
        with open(temp_file_path, "wb") as temp_file:
            temp_file.write(uploaded_file.read())

        # Display PDF and summarize on button click
        if st.button("Summarize Document"):
            col1, col2 = st.columns(2)
            with col1:
                st.info("Uploaded File")
                if file_type == "pdf":
                    displayPDF(temp_file_path)

            with col2:
                summary = llm_pipeline(temp_file_path, file_type)
                st.info("Summarization Complete")
                st.success(summary)

    # Image description section
    st.header("Image Description")
    uploaded_image = st.file_uploader("Upload an Image (PNG, JPG, JPEG)", type=['png', 'jpg', 'jpeg'])

    if uploaded_image is not None:
        image = Image.open(uploaded_image)
        st.image(image, caption="Uploaded Image", use_column_width=True)

        if st.button("Describe Image"):
            description = describe_image(uploaded_image)
            st.info("Image Description")
            st.success(description)

if __name__ == "__main__":
    main()